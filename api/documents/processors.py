import io
import aiofiles
from pyuca import Collator
import pandas as pd
import traceback
from pathlib import Path
from typing import Dict, List, Optional, Any, Literal
from lightrag import LightRAG
from lightrag.utils import logger

# Temporary file prefix
temp_prefix = "__tmp__"

class DocumentManager:
    def __init__(
        self,
        input_dir: str,
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".rtf",  # Rich Text Format
            ".odt",  # OpenDocument Text
            ".tex",  # LaTeX
            ".epub",  # Electronic Publication
            ".html",  # HyperText Markup Language
            ".htm",  # HyperText Markup Language
            ".csv",  # Comma-Separated Values
            ".json",  # JavaScript Object Notation
            ".xml",  # eXtensible Markup Language
            ".yaml",  # YAML Ain't Markup Language
            ".yml",  # YAML
            ".log",  # Log files
            ".conf",  # Configuration files
            ".ini",  # Initialization files
            ".properties",  # Java properties files
            ".sql",  # SQL scripts
            ".bat",  # Batch files
            ".sh",  # Shell scripts
            ".c",  # C source code
            ".cpp",  # C++ source code
            ".py",  # Python source code
            ".java",  # Java source code
            ".js",  # JavaScript source code
            ".ts",  # TypeScript source code
            ".swift",  # Swift source code
            ".go",  # Go source code
            ".rb",  # Ruby source code
            ".php",  # PHP source code
            ".css",  # Cascading Style Sheets
            ".scss",  # Sassy CSS
            ".less",  # LESS CSS
        ),
    ):
        self.input_dir = Path(input_dir)
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            logger.debug(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.rglob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


async def pipeline_enqueue_file(rag: LightRAG, doc_id:str, file_path: Path) -> bool:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        doc_id: Identifier for the document
        file_path: Path to the saved file
    Returns:
        bool: True if the file was successfully enqueued, False otherwise
    """

    try:
        content = ""
        ext = file_path.suffix.lower()

        file = None
        async with aiofiles.open(file_path, "rb") as f:
            file = await f.read()

        # Process based on file type
        match ext:
            case (
                ".txt"
                | ".md"
                | ".html"
                | ".htm"
                | ".tex"
                | ".json"
                | ".xml"
                | ".yaml"
                | ".yml"
                | ".rtf"
                | ".odt"
                | ".epub"
                | ".log"
                | ".conf"
                | ".ini"
                | ".properties"
                | ".sql"
                | ".bat"
                | ".sh"
                | ".c"
                | ".cpp"
                | ".py"
                | ".java"
                | ".js"
                | ".ts"
                | ".swift"
                | ".go"
                | ".rb"
                | ".php"
                | ".css"
                | ".scss"
                | ".less"
            ):
                try:
                    # Try to decode as UTF-8
                    content = file.decode("utf-8")

                    # Validate content
                    if not content or len(content.strip()) == 0:
                        logger.error(f"Empty content in file: {file_path.name}")
                        return False

                    # Check if content looks like binary data string representation
                    if content.startswith("b'") or content.startswith('b"'):
                        logger.error(
                            f"File {file_path.name} appears to contain binary data representation instead of text"
                        )
                        return False

                except UnicodeDecodeError:
                    logger.error(
                        f"File {file_path.name} is not valid UTF-8 encoded text. Please convert it to UTF-8 before processing."
                    )
                    return False
            case ".pdf":
                    from PyPDF2 import PdfReader  # type: ignore
                    from io import BytesIO

                    pdf_file = BytesIO(file)
                    reader = PdfReader(pdf_file)
                    for i, page in enumerate(reader.pages):
                        content += f"Page {i + 1}:\n\n " + page.extract_text() + "\n"

            case ".docx" | ".doc":
                try:
                    from docling.document_converter import DocumentConverter  # type: ignore
                    converter = DocumentConverter()
                    result = converter.convert(file_path)
                    content = result.document.export_to_markdown()
                except:
                    from docx import Document  # type: ignore
                    from io import BytesIO

                    docx_file = BytesIO(file)
                    doc = Document(docx_file)
                    content = "\n".join(
                        [paragraph.text for paragraph in doc.paragraphs]
                    )
            case ".pptx":
                try:
                    from pptx import Presentation  # type: ignore
                    from io import BytesIO

                    pptx_file = BytesIO(file)
                    prs = Presentation(pptx_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                except:
                    logger.error(
                        f"Error processing PPTX file {file_path.name}. Ensure python-pptx is installed."
                    )

            case '.csv':
                try:
                    df = pd.read_csv(io.BytesIO(file))
                    sheet_text = f"\n\n--- Filename: {file_path.name} ---\n\n"
                    records = df.to_dict(orient='records')
                    content +=  sheet_text + str(records)
                    logger.debug(f"CSV processing content {content[:100]}......")
                except Exception as e:
                    logger.error(f"Error processing CSV file: {str(e)}")

            case '.xlsx' | '.xls' | '.xlsm':
                try:
                    excel_file = io.BytesIO(file)
                    all_sheets_content = ""
                    excel_data = pd.read_excel(excel_file, sheet_name=None)

                    logger.info(f"Excel file read, found {len(excel_data)} sheets")
                    
                    # Process each sheet
                    num_sheets = 0
                    for sheet_name, df in excel_data.items():
                        logger.debug(f"Processing sheet: {sheet_name} with {len(df)} rows")
                        sheet_text = f"\n\n--- Sheet: {sheet_name} ---\n\n"
                        records = df.to_dict(orient='records')
                        all_sheets_content += f"{sheet_text}  {str(records)}"
                        num_sheets += 1
                        
                    # Join all sheets with separator
                    content = all_sheets_content
                    logger.info(f"Processed Excel file with {num_sheets} sheets")
                except Exception as e:
                    logger.error(f"Error processing Excel file: {str(e)}")
                    
            case _:
                logger.error(
                    f"Unsupported file type: {file_path.name} (extension {ext})"
                )
                return False

        # Insert into the RAG queue
        if content:
            await rag.apipeline_enqueue_documents(content, ids=[doc_id], file_paths=file_path.name)
            logger.info(f"Successfully fetched and enqueued file: {file_path.name}")
            return True
        else:
            logger.error(f"No content could be extracted from file: {file_path.name}")

    except Exception as e:
        logger.error(f"Error processing or enqueueing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())
    finally:
        if file_path.name.startswith(temp_prefix):
            try:
                file_path.unlink()
            except Exception as e:
                logger.error(f"Error deleting file {file_path}: {str(e)}")
    return False


async def pipeline_index_file(rag: LightRAG, doc_id: str, file_path: Path):
    """Index a file

    Args:
        rag: LightRAG instance
        doc_id: Document ID for the file
        file_path: Path to the saved file
    """
    try:
        if await pipeline_enqueue_file(rag, doc_id, file_path):
            await rag.apipeline_process_enqueue_documents()

    except Exception as e:
        logger.error(f"Error indexing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_files(rag: LightRAG, doc_ids: List[str], file_paths: List[Path]):
    """Index multiple files sequentially to avoid high CPU load

    Args:
        rag: LightRAG instance
        doc_ids: Document IDs for the files
        file_paths: Paths to the files to index
    """
    if not file_paths:
        return
    try:
        enqueued = False

        # Create Collator for Unicode sorting
        collator = Collator()
        sorted_file_paths = sorted(file_paths, key=lambda p: collator.sort_key(str(p)))

        # Process files sequentially
        for doc_id, file_path in zip(doc_ids, sorted_file_paths):
            if await pipeline_enqueue_file(rag, doc_id, file_path):
                enqueued = True

        # Process the queue only if at least one file was successfully enqueued
        if enqueued:
            await rag.apipeline_process_enqueue_documents()
    except Exception as e:
        logger.error(f"Error indexing files: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_texts(
    rag: LightRAG, texts: List[str], file_sources: List[str] = None
):
    """Index a list of texts

    Args:
        rag: LightRAG instance
        texts: The texts to index
        file_sources: Sources of the texts
    """
    if not texts:
        return
    if file_sources is not None:
        if len(file_sources) != 0 and len(file_sources) != len(texts):
            [
                file_sources.append("unknown_source")
                for _ in range(len(file_sources), len(texts))
            ]
    await rag.apipeline_enqueue_documents(input=texts, file_paths=file_sources)
    await rag.apipeline_process_enqueue_documents()




